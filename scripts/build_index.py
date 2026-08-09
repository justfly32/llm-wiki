#!/usr/bin/env python3
"""
LLM Wiki — Build Index (v2, Glean 형태)
==========================================
Hermes Agent가 만들어낸 모든 작업 결과물을 SQLite FTS5 인덱스로 구축한다.

인덱싱 대상:
  - ~/projects/   (모든 프로젝트: 코드, 문서, 설정)
  - ~/.hermes/    (Hermes: scripts, skills, memories, cron, kanban 등)
  - ~/documents/  (작업 문서)

동작:
  1. 루트별 재귀 스캔 (노이즈 디렉토리 제외)
  2. 파일 유형 분류 → 텍스트 추출 (코드/마크다운/JSON/설정/문서/PDF)
  3. SQLite FTS5 인덱스에 저장 (증분: mtime 비교)

사용법:
  python3 build_index.py            # 증분 인덱싱
  python3 build_index.py --full     # 전체 재인덱싱
  python3 build_index.py --stats    # 인덱스 통계
"""
import os
import re
import sys
import json
import sqlite3
import argparse
from datetime import datetime
from pathlib import Path

HOME = Path.home()
INDEX_DB = Path.home() / "wiki" / "index.db"
STATE_FILE = Path("/tmp/.llm-wiki-index-state")

SCAN_ROOTS = [
    ("projects", HOME / "projects"),
    ("hermes", HOME / ".hermes"),
    ("documents", HOME / "documents"),
    ("wiki", HOME / "wiki"),  # 위키 페이지(concepts/)도 인덱싱 — 중복 방지 문서 검색용 (2026-08-09)
]

EXCLUDE_DIRS = {
    "node_modules", ".git", ".venv", "venv", "__pycache__", "dist", "build",
    ".next", ".cache", ".pytest_cache", ".mypy_cache", ".ruff_cache",
    "logs", "cache", "audio_cache", "image_cache", "bootstrap-cache", "node",
    "lsp", "pastes", "sandboxes", "sessions", "rate_limits", "migration",
    "pairing", "profiles", "gateway", "mcp-env", "bin", "shared",
    "state-snapshots", "venvs", "output", "bootstrap-cache",
    ".venv", "site-packages", "Pods", ".build", "DerivedData", ".gradle",
}

EXCLUDE_FILES = {
    "state.db", "state.db-wal", "state.db-shm", "verification_evidence.db",
    "models_dev_cache.json", "gateway_state.json", "channel_directory.json",
    "ticker_heartbeat", "ticker_last_success", ".tick.lock", ".jobs.lock",
    "index.db", "package-lock.json", "pnpm-lock.yaml", "yarn.lock",
    ".DS_Store", "LICENSE",
}

# 텍스트 추출 대상 확장자 (바이너리 제외)
TEXT_EXTS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".sh", ".bash", ".go", ".rs",
    ".java", ".c", ".cpp", ".h", ".swift", ".kt", ".rb", ".php", ".sql",
    ".md", ".mdx", ".txt", ".rst", ".json", ".yaml", ".yml", ".toml",
    ".ini", ".env", ".conf", ".plist", ".cfg", ".log", ".csv", ".html",
    ".css", ".scss", ".svg", ".xml", ".vue", ".svelte", ".ipynb",
}
DOC_EXTS = {".pdf", ".docx", ".hwp", ".xlsx", ".pptx"}

TYPE_MAP = {
    ".py": "code", ".js": "code", ".ts": "code", ".tsx": "code", ".jsx": "code",
    ".sh": "code", ".bash": "code", ".go": "code", ".rs": "code", ".java": "code",
    ".c": "code", ".cpp": "code", ".h": "code", ".swift": "code", ".kt": "code",
    ".rb": "code", ".php": "code", ".sql": "code", ".vue": "code", ".svelte": "code",
    ".md": "doc", ".mdx": "doc", ".txt": "doc", ".rst": "doc", ".ipynb": "doc",
    ".json": "config", ".yaml": "config", ".yml": "config", ".toml": "config",
    ".ini": "config", ".env": "config", ".conf": "config", ".plist": "config",
    ".html": "web", ".css": "web", ".scss": "web", ".svg": "image",
    ".csv": "data", ".xml": "data",
    ".pdf": "pdf", ".docx": "docx", ".hwp": "hwp", ".xlsx": "xlsx", ".pptx": "pptx",
}

def classify(path: Path) -> str:
    ext = path.suffix.lower()
    return TYPE_MAP.get(ext, "other")

def should_skip_dir(name: str) -> bool:
    return name in EXCLUDE_DIRS or name.startswith(".")

def extract_text(path: Path, ftype: str) -> str:
    """파일에서 텍스트 추출 (크기 제한 512KB)"""
    try:
        size = path.stat().st_size
        if size > 512 * 1024:
            return ""
        if ftype in ("pdf", "docx", "hwp", "xlsx", "pptx"):
            return extract_doc(path, ftype)
        raw = path.read_bytes()
        # 바이너리 감지 (NUL 바이트 비율)
        if b"\x00" in raw[:2048]:
            return ""
        return raw.decode("utf-8", errors="replace")
    except (OSError, PermissionError):
        return ""

def extract_doc(path: Path, ftype: str) -> str:
    """문서 파일 텍스트 추출 (선택적 의존성)"""
    try:
        if ftype == "pdf":
            import fitz  # PyMuPDF
            doc = fitz.open(path)
            return "\n".join(page.get_text() for page in doc)[:200_000]
        elif ftype == "docx":
            import docx
            d = docx.Document(path)
            return "\n".join(p.text for p in d.paragraphs)[:200_000]
        elif ftype == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append(" | ".join(str(c) for c in row if c is not None))
            return "\n".join(parts)[:200_000]
        elif ftype == "pptx":
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)[:200_000]
        elif ftype == "hwp":
            # HWP는 복잡 — 기본적으로 스킵 (메타데이터만)
            return ""
    except ImportError:
        return ""
    except Exception:
        return ""

def connect_db():
    conn = sqlite3.connect(INDEX_DB)
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("""
        CREATE TABLE IF NOT EXISTS files (
            path TEXT PRIMARY KEY,
            rel TEXT,
            root TEXT,
            type TEXT,
            size INTEGER,
            mtime REAL,
            content TEXT
        )
    """)
    conn.execute("""
        CREATE VIRTUAL TABLE IF NOT EXISTS files_fts USING fts5(
            path, rel, root, type, content, tokenize='unicode61'
        )
    """)
    return conn

def sync_fts(conn):
    """files 테이블 변경분을 FTS로 동기화 (간단 전체 재구축)"""
    conn.execute("DELETE FROM files_fts")
    conn.execute("""
        INSERT INTO files_fts(path, rel, root, type, content)
        SELECT path, rel, root, type, content FROM files WHERE content != ''
    """)

def scan_and_index(conn, full=False):
    """루트 스캔 + 인덱싱"""
    # 증분 기준 시각
    since = 0.0
    if not full and STATE_FILE.exists():
        try:
            since = float(STATE_FILE.read_text().strip())
        except ValueError:
            since = 0.0

    indexed = 0
    skipped = 0
    new_files = 0
    updated = 0

    for root_name, root in SCAN_ROOTS:
        if not root.exists():
            continue
        for dirpath, dirnames, filenames in os.walk(root):
            dirnames[:] = [d for d in dirnames if not should_skip_dir(d)]
            for fname in filenames:
                if fname in EXCLUDE_FILES:
                    continue
                fp = Path(dirpath) / fname
                ext = fp.suffix.lower()
                if ext not in TEXT_EXTS and ext not in DOC_EXTS:
                    continue
                try:
                    st = fp.stat()
                except OSError:
                    continue
                # 증분: mtime이 기준보다 오래된 파일은 스킵 (full이면 전부)
                if not full and st.st_mtime <= since:
                    continue
                rel = str(fp.relative_to(root))
                ftype = classify(fp)
                content = extract_text(fp, ftype)
                if not content.strip():
                    skipped += 1

                # 기존 여부 확인
                cur = conn.execute("SELECT mtime FROM files WHERE path=?", (str(fp),)).fetchone()
                if cur is not None:
                    if not full and abs(cur[0] - st.st_mtime) < 1:
                        continue  # 변경 없음
                    conn.execute(
                        "UPDATE files SET rel=?, root=?, type=?, size=?, mtime=?, content=? WHERE path=?",
                        (rel, root_name, ftype, st.st_size, st.st_mtime, content, str(fp)))
                    updated += 1
                else:
                    conn.execute(
                        "INSERT INTO files(path, rel, root, type, size, mtime, content) VALUES (?,?,?,?,?,?,?)",
                        (str(fp), rel, root_name, ftype, st.st_size, st.st_mtime, content))
                    new_files += 1
                indexed += 1

    conn.commit()
    sync_fts(conn)
    conn.commit()

    # 상태 기록
    STATE_FILE.write_text(str(datetime.now().timestamp()))

    return {"indexed": indexed, "new": new_files, "updated": updated, "skipped": skipped}

def print_stats(conn):
    total = conn.execute("SELECT COUNT(*) FROM files").fetchone()[0]
    fts = conn.execute("SELECT COUNT(*) FROM files_fts").fetchone()[0]
    by_root = conn.execute("SELECT root, COUNT(*) FROM files GROUP BY root").fetchall()
    by_type = conn.execute("SELECT type, COUNT(*) FROM files GROUP BY type ORDER BY 2 DESC").fetchall()
    size = conn.execute("SELECT SUM(size) FROM files").fetchone()[0] or 0
    print(f"📊 인덱스 통계: 총 {total}개 파일 (FTS {fts}개, {size/1024/1024:.1f}MB)")
    print(f"  루트별: {', '.join(f'{r} {c}개' for r, c in by_root)}")
    print(f"  유형별: {', '.join(f'{t} {c}개' for t, c in by_type[:8])}")

def main():
    parser = argparse.ArgumentParser(description="LLM Wiki 인덱스 빌더")
    parser.add_argument("--full", action="store_true", help="전체 재인덱싱")
    parser.add_argument("--stats", action="store_true", help="통계만 출력")
    args = parser.parse_args()

    conn = connect_db()
    if args.stats:
        print_stats(conn)
        return

    print("🔄 작업 결과물 인덱싱 중...")
    start = datetime.now()
    result = scan_and_index(conn, full=args.full)
    elapsed = (datetime.now() - start).total_seconds()

    print(f"✅ 인덱싱 완료 ({elapsed:.1f}초)")
    print(f"  신규 {result['new']}개 / 갱신 {result['updated']}개 / 텍스트 없음 스킵 {result['skipped']}개")
    print_stats(conn)

if __name__ == "__main__":
    main()
